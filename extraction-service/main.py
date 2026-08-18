from fastapi import FastAPI, UploadFile, File, HTTPException
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
from PIL import Image
from pdf2image import convert_from_path
from striprtf.striprtf import rtf_to_text

import pandas as pd
import pytesseract
import tempfile
import subprocess
import shutil
import os
import json
import csv
import xml.etree.ElementTree as ET

from ebooklib import epub
from odf import text, teletype
from odf.opendocument import load as odf_load


app = FastAPI(
    title="ASP Document Extraction Service",
    version="2.0.0",
    description="Universal document extraction service for ASP"
)


# ---------------------------------------------------------
# HEALTH
# ---------------------------------------------------------

@app.get("/health")
def health():
    return {
        "status": "ok",
        "service": "asp-extraction-service",
        "version": "2.0.0"
    }


# ---------------------------------------------------------
# UTILITY
# ---------------------------------------------------------

def clean_text(value):
    if value is None:
        return ""

    return " ".join(str(value).replace("\x00", " ").split())


def make_element(element_type, text_content, page=None, filename=None):
    text_content = clean_text(text_content)

    if not text_content:
        return None

    return {
        "type": element_type,
        "text": text_content,
        "metadata": {
            "page_number": page,
            "filename": filename
        }
    }


# ---------------------------------------------------------
# PDF
# ---------------------------------------------------------

def extract_pdf(path, filename):
    elements = []

    reader = PdfReader(path)

    for page_number, page in enumerate(reader.pages, start=1):

        try:
            text_content = page.extract_text() or ""
        except Exception:
            text_content = ""

        text_content = clean_text(text_content)

        if text_content:
            element = make_element(
                "Page",
                text_content,
                page_number,
                filename
            )

            if element:
                elements.append(element)

    # OCR fallback for scanned PDFs
    if not elements:
        try:
            images = convert_from_path(
                path,
                dpi=150,
                fmt="jpeg"
            )

            for page_number, image in enumerate(images, start=1):

                try:
                    ocr_text = pytesseract.image_to_string(image)
                except Exception:
                    ocr_text = ""

                element = make_element(
                    "OCRPage",
                    ocr_text,
                    page_number,
                    filename
                )

                if element:
                    elements.append(element)

        except Exception as error:
            raise RuntimeError(
                f"PDF OCR failed: {str(error)}"
            )

    return elements


# ---------------------------------------------------------
# DOCX
# ---------------------------------------------------------

def extract_docx(path, filename):
    elements = []

    document = Document(path)

    for paragraph in document.paragraphs:

        element = make_element(
            "Paragraph",
            paragraph.text,
            None,
            filename
        )

        if element:
            elements.append(element)

    for table_index, table in enumerate(document.tables, start=1):

        rows = []

        for row in table.rows:
            cells = [
                clean_text(cell.text)
                for cell in row.cells
            ]

            rows.append(" | ".join(cells))

        table_text = "\n".join(rows)

        element = make_element(
            "Table",
            table_text,
            None,
            filename
        )

        if element:
            element["metadata"]["table_number"] = table_index
            elements.append(element)

    return elements


# ---------------------------------------------------------
# PPTX
# ---------------------------------------------------------

def extract_pptx(path, filename):
    elements = []

    presentation = Presentation(path)

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text_content = clean_text(shape.text)

                if text_content:
                    slide_text.append(text_content)

        element = make_element(
            "Slide",
            "\n".join(slide_text),
            slide_number,
            filename
        )

        if element:
            elements.append(element)

    return elements


# ---------------------------------------------------------
# XLSX / XLS
# ---------------------------------------------------------

def extract_excel(path, filename):
    elements = []

    extension = os.path.splitext(path)[1].lower()

    if extension == ".xlsx":

        workbook = load_workbook(
            path,
            read_only=True,
            data_only=True
        )

        for sheet in workbook.worksheets:

            rows = []

            for row in sheet.iter_rows(values_only=True):

                values = [
                    clean_text(value)
                    for value in row
                ]

                if any(values):
                    rows.append(" | ".join(values))

            sheet_text = "\n".join(rows)

            element = make_element(
                "Spreadsheet",
                sheet_text,
                None,
                filename
            )

            if element:
                element["metadata"]["sheet"] = sheet.title
                elements.append(element)

    else:

        excel_file = pd.ExcelFile(path)

        for sheet_name in excel_file.sheet_names:

            dataframe = pd.read_excel(
                path,
                sheet_name=sheet_name,
                header=None
            )

            rows = []

            for row in dataframe.itertuples(
                index=False,
                name=None
            ):

                values = [
                    clean_text(value)
                    for value in row
                ]

                if any(values):
                    rows.append(" | ".join(values))

            element = make_element(
                "Spreadsheet",
                "\n".join(rows),
                None,
                filename
            )

            if element:
                element["metadata"]["sheet"] = sheet_name
                elements.append(element)

    return elements


# ---------------------------------------------------------
# CSV
# ---------------------------------------------------------

def extract_csv(path, filename):
    elements = []

    with open(
        path,
        "r",
        encoding="utf-8",
        errors="replace",
        newline=""
    ) as file:

        reader = csv.reader(file)

        rows = []

        for row in reader:
            values = [
                clean_text(value)
                for value in row
            ]

            rows.append(" | ".join(values))

    element = make_element(
        "CSV",
        "\n".join(rows),
        None,
        filename
    )

    if element:
        elements.append(element)

    return elements


# ---------------------------------------------------------
# TXT / MD
# ---------------------------------------------------------

def extract_text_file(path, filename):
    with open(
        path,
        "r",
        encoding="utf-8",
        errors="replace"
    ) as file:

        content = file.read()

    element = make_element(
        "Text",
        content,
        None,
        filename
    )

    return [element] if element else []


# ---------------------------------------------------------
# HTML
# ---------------------------------------------------------

def extract_html(path, filename):
    with open(
        path,
        "r",
        encoding="utf-8",
        errors="replace"
    ) as file:

        html = file.read()

    soup = BeautifulSoup(
        html,
        "lxml"
    )

    for tag in soup([
        "script",
        "style",
        "noscript"
    ]):
        tag.decompose()

    content = soup.get_text(
        separator="\n"
    )

    element = make_element(
        "HTML",
        content,
        None,
        filename
    )

    return [element] if element else []


# ---------------------------------------------------------
# XML
# ---------------------------------------------------------

def extract_xml(path, filename):
    tree = ET.parse(path)

    root = tree.getroot()

    content = " ".join(
        root.itertext()
    )

    element = make_element(
        "XML",
        content,
        None,
        filename
    )

    return [element] if element else []


# ---------------------------------------------------------
# JSON
# ---------------------------------------------------------

def extract_json(path, filename):
    with open(
        path,
        "r",
        encoding="utf-8",
        errors="replace"
    ) as file:

        data = json.load(file)

    content = json.dumps(
        data,
        indent=2,
        ensure_ascii=False
    )

    element = make_element(
        "JSON",
        content,
        None,
        filename
    )

    return [element] if element else []


# ---------------------------------------------------------
# IMAGE OCR
# ---------------------------------------------------------

def extract_image(path, filename):
    image = Image.open(path)

    content = pytesseract.image_to_string(
        image
    )

    element = make_element(
        "OCR",
        content,
        None,
        filename
    )

    return [element] if element else []


# ---------------------------------------------------------
# RTF
# ---------------------------------------------------------

def extract_rtf(path, filename):
    with open(
        path,
        "r",
        encoding="utf-8",
        errors="replace"
    ) as file:

        content = file.read()

    text_content = rtf_to_text(
        content
    )

    element = make_element(
        "RTF",
        text_content,
        None,
        filename
    )

    return [element] if element else []


# ---------------------------------------------------------
# EPUB
# ---------------------------------------------------------

def extract_epub(path, filename):
    elements = []

    book = epub.read_epub(path)

    for item in book.get_items():

        if item.get_type() == 9:

            soup = BeautifulSoup(
                item.get_content(),
                "lxml"
            )

            content = soup.get_text(
                separator="\n"
            )

            element = make_element(
                "EPUB",
                content,
                None,
                filename
            )

            if element:
                elements.append(element)

    return elements


# ---------------------------------------------------------
# ODT / ODS / ODP
# ---------------------------------------------------------

def extract_odf(path, filename):
    elements = []

    document = odf_load(path)

    paragraphs = document.getElementsByType(
        text.P
    )

    for paragraph in paragraphs:

        content = teletype.extractText(
            paragraph
        )

        element = make_element(
            "ODF",
            content,
            None,
            filename
        )

        if element:
            elements.append(element)

    return elements


# ---------------------------------------------------------
# LEGACY OFFICE FILES
# DOC / PPT
# ---------------------------------------------------------

def convert_with_libreoffice(
    path,
    output_dir
):

    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "docx",
        "--outdir",
        output_dir,
        path
    ]

    subprocess.run(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=120
    )

    converted = os.path.join(
        output_dir,
        os.path.splitext(
            os.path.basename(path)
        )[0] + ".docx"
    )

    return converted if os.path.exists(
        converted
    ) else None


# ---------------------------------------------------------
# UNIVERSAL DISPATCHER
# ---------------------------------------------------------

def extract_document(
    path,
    filename
):

    extension = os.path.splitext(
        filename
    )[1].lower()

    if extension == ".pdf":
        return extract_pdf(path, filename)

    if extension == ".docx":
        return extract_docx(path, filename)

    if extension == ".pptx":
        return extract_pptx(path, filename)

    if extension in [".xlsx", ".xls"]:
        return extract_excel(path, filename)

    if extension == ".csv":
        return extract_csv(path, filename)

    if extension in [
        ".txt",
        ".md",
        ".markdown"
    ]:
        return extract_text_file(
            path,
            filename
        )

    if extension in [
        ".html",
        ".htm"
    ]:
        return extract_html(
            path,
            filename
        )

    if extension == ".xml":
        return extract_xml(
            path,
            filename
        )

    if extension == ".json":
        return extract_json(
            path,
            filename
        )

    if extension in [
        ".png",
        ".jpg",
        ".jpeg",
        ".tiff",
        ".tif",
        ".bmp",
        ".webp"
    ]:
        return extract_image(
            path,
            filename
        )

    if extension == ".rtf":
        return extract_rtf(
            path,
            filename
        )

    if extension == ".epub":
        return extract_epub(
            path,
            filename
        )

    if extension in [
        ".odt",
        ".ods",
        ".odp"
    ]:
        return extract_odf(
            path,
            filename
        )

    if extension in [
        ".doc",
        ".ppt"
    ]:

        temp_dir = tempfile.mkdtemp()

        try:

            converted = convert_with_libreoffice(
                path,
                temp_dir
            )

            if not converted:
                raise RuntimeError(
                    "LibreOffice could not convert "
                    f"{extension} file"
                )

            if extension == ".doc":
                return extract_docx(
                    converted,
                    filename
                )

            return extract_pptx(
                converted,
                filename
            )

        finally:
            shutil.rmtree(
                temp_dir,
                ignore_errors=True
            )

    raise ValueError(
        f"Unsupported document format: {extension}"
    )


# ---------------------------------------------------------
# EXTRACT ENDPOINT
# ---------------------------------------------------------

@app.post("/extract")
async def extract_document_endpoint(
    file: UploadFile = File(...)
):

    if not file.filename:
        raise HTTPException(
            status_code=400,
            detail="No filename provided"
        )

    temp_path = None

    try:

        suffix = os.path.splitext(
            file.filename
        )[1]

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=suffix
        ) as temp_file:

            contents = await file.read()

            temp_file.write(
                contents
            )

            temp_path = temp_file.name

        extracted_elements = extract_document(
            temp_path,
            file.filename
        )

        full_text = "\n\n".join(
            item["text"]
            for item in extracted_elements
            if item
        )

        return {
            "success": True,
            "filename": file.filename,
            "file_type": suffix.lower(),
            "element_count": len(
                extracted_elements
            ),
            "text": full_text,
            "elements": extracted_elements
        }

    except ValueError as error:

        raise HTTPException(
            status_code=415,
            detail=str(error)
        )

    except Exception as error:

        raise HTTPException(
            status_code=500,
            detail=(
                "Document extraction failed: "
                f"{str(error)}"
            )
        )

    finally:

        if (
            temp_path
            and os.path.exists(temp_path)
        ):
            os.remove(temp_path)